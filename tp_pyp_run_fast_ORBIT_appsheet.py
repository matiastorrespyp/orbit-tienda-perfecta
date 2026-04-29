#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import argparse
import json
import os
import re
import difflib
import unicodedata
from collections import Counter
from datetime import date, datetime, timedelta
from math import ceil

import numpy as np
import pandas as pd

from reportlab.lib.pagesizes import A4, landscape
from reportlab.pdfgen import canvas
from reportlab.lib.units import cm
from reportlab.lib import colors


# ============================================================
# Firma / Footer
# ============================================================
CREATOR_NAME = "Matías Torres"
SCRIPT_VERSION = "ORBIT-STABLE-2026-04-28"


def orbit_logo_path():
    return os.path.join(script_dir(), "assets", "orbit.png")

def draw_footer_signature(c: canvas.Canvas, page_w: float, page_h: float):
    c.saveState()
    footer_y = 0.58 * cm
    logo = orbit_logo_path()
    logo_w = 2.1 * cm
    logo_h = 0.78 * cm
    text_x = 1.55 * cm
    if os.path.exists(logo):
        try:
            c.drawImage(
                logo,
                1.0 * cm,
                footer_y - 0.14 * cm,
                width=logo_w,
                height=logo_h,
                preserveAspectRatio=True,
                mask="auto",
            )
            text_x = 1.0 * cm + logo_w + 0.18 * cm
        except Exception:
            text_x = 1.55 * cm
    try:
        c.setFillAlpha(0.42)
    except Exception:
        pass
    c.setFont("Helvetica", 8.2)
    c.setFillColor(colors.Color(0, 0, 0, alpha=0.42))
    c.drawString(text_x, footer_y + 0.02 * cm, f"Creado por {CREATOR_NAME}")
    c.drawRightString(page_w - 1.2 * cm, footer_y + 0.02 * cm, f"Pág. {c.getPageNumber()}")
    c.restoreState()


def finalize_page(c: canvas.Canvas, page_w: float, page_h: float):
    draw_footer_signature(c, page_w, page_h)


# ============================================================
# Helpers
# ============================================================
def norm_text(x):
    if x is None or (isinstance(x, float) and np.isnan(x)):
        return ""
    s = str(x).strip()
    s = re.sub(r"\s+", " ", s)
    return s


def norm_vendor_id(x):
    s = norm_text(x)
    if not s:
        return ""
    m = re.fullmatch(r"(\d+)\.0+", s)
    if m:
        return m.group(1)
    return s

def safe_mkdir(p):
    os.makedirs(p, exist_ok=True)

def to_num(x):
    return pd.to_numeric(x, errors="coerce")

def to_pct_0_100(series: pd.Series) -> pd.Series:
    s = pd.to_numeric(series, errors="coerce")
    if s.notna().sum() == 0:
        return s
    q95 = s.quantile(0.95)
    if pd.notna(q95) and q95 <= 1.5:
        s = s * 100.0
    return s

def parse_listish(cell):
    if cell is None or (isinstance(cell, float) and np.isnan(cell)):
        return []
    if isinstance(cell, list):
        return [norm_text(x) for x in cell if norm_text(x)]
    s = norm_text(cell)
    if not s:
        return []
    parts = re.split(r"[;|]\s*|,\s*", s)
    return [p for p in (norm_text(x) for x in parts) if p]

def script_dir():
    return os.path.dirname(os.path.abspath(__file__))


def resolve_input_file(path_value, patterns):
    candidates = []
    if path_value:
        candidates.append(path_value)
        if not os.path.isabs(path_value):
            candidates.append(os.path.join(script_dir(), path_value))
            candidates.append(os.path.join(script_dir(), "inputs", path_value))
    for cand in candidates:
        if cand and os.path.exists(cand):
            return cand

    search_dirs = [script_dir(), os.path.join(script_dir(), "inputs")]
    for base in search_dirs:
        if not os.path.isdir(base):
            continue
        for name in sorted(os.listdir(base)):
            low = name.lower()
            full = os.path.join(base, name)
            if not os.path.isfile(full):
                continue
            for patt in patterns:
                if re.search(patt, low):
                    return full
    return path_value


def day_code_for_tomorrow(d=None):
    if d is None:
        d = date.today()
    t = d + timedelta(days=1)
    wd = t.weekday()
    m = {0: "LU", 1: "MA", 2: "MI", 3: "JU", 4: "VI", 5: "SA", 6: "DO"}
    code = m[wd]
    return "LU" if code == "DO" else code

def freq_has_day(freq_val, dia_code):
    if freq_val is None or (isinstance(freq_val, float) and np.isnan(freq_val)):
        return False
    s = norm_text(freq_val).upper()
    if not s:
        return False
    tokens = re.findall(r"[A-Z]{2}", s)
    return dia_code in tokens

def fmt_pesos(n):
    try:
        return "$" + f"{float(n):,.0f}".replace(",", ".")
    except Exception:
        return "$0"



def parse_decimal_series(series: pd.Series) -> pd.Series:
    s = series.astype(str).str.strip()
    s = s.replace({"": np.nan, "nan": np.nan, "None": np.nan, "NULL": np.nan})
    s = s.str.replace(r"[$\s]", "", regex=True)
    both = s.str.contains(",", na=False) & s.str.contains(r"\.", na=False)
    s.loc[both] = s.loc[both].str.replace(".", "", regex=False).str.replace(",", ".", regex=False)
    only_comma = s.str.contains(",", na=False) & ~s.str.contains(r"\.", na=False)
    s.loc[only_comma] = s.loc[only_comma].str.replace(",", ".", regex=False)
    return pd.to_numeric(s, errors="coerce")

def detect_sales_amount_column(df: pd.DataFrame):
    priority = [
        "ImporteNetoItem", "ImporteItem", "NetoItem",
        "Importe Neto", "ImporteNeto", "Neto",
        "Total sin IVA", "TotalSinIVA", "Subtotal sin IVA",
        "Importe", "Total",
    ]
    for c in priority:
        if c in df.columns:
            vals = parse_decimal_series(df[c]) if df[c].dtype == object else pd.to_numeric(df[c], errors="coerce")
            if vals.notna().sum() > 0 and float(vals.abs().sum()) > 0:
                return c
    best_col, best_score = None, -1e18
    for c in df.columns:
        vals = parse_decimal_series(df[c]) if df[c].dtype == object else pd.to_numeric(df[c], errors="coerce")
        valid = vals.notna().sum()
        if valid == 0:
            continue
        ck = norm_key(c)
        score = float(valid)
        if "IMPORTE" in ck:
            score += 300
        if "NETO" in ck:
            score += 250
        if "TOTAL" in ck:
            score += 180
        if "SUBTOTAL" in ck:
            score += 140
        if "CANT" in ck or "BASE" in ck or "PRECIO" in ck or "COSTO" in ck:
            score -= 500
        if score > best_score and float(vals.abs().sum()) > 0:
            best_score = score
            best_col = c
    return best_col


def detect_sales_weight_column(df: pd.DataFrame):
    priority = [
        "PesoKg", "PesoKG", "Kilos", "Kg", "Kilogramos", "Peso",
        "PesoTotalKg", "CantidadKg",
    ]
    for c in priority:
        if c in df.columns:
            vals = parse_decimal_series(df[c]) if df[c].dtype == object else pd.to_numeric(df[c], errors="coerce")
            if vals.notna().sum() > 0 and float(vals.abs().sum()) > 0:
                return c
    best_col, best_score = None, -1e18
    for c in df.columns:
        vals = parse_decimal_series(df[c]) if df[c].dtype == object else pd.to_numeric(df[c], errors="coerce")
        valid = vals.notna().sum()
        if valid == 0:
            continue
        ck = norm_key(c)
        score = float(valid)
        if "PESO" in ck:
            score += 350
        if ck in ("KG", "KILOS", "KILOGRAMOS"):
            score += 320
        if "CANT" in ck or "PRECIO" in ck or "IMPORTE" in ck or "NETO" in ck or "TOTAL" in ck:
            score -= 400
        if score > best_score and float(vals.abs().sum()) > 0:
            best_score = score
            best_col = c
    return best_col

def attach_sales_metrics_from_sales(gescom_df, sales_df):
    g = gescom_df.copy()
    g["FACTURACION_CLIENTE"] = 0.0
    g["KILOS_CLIENTE"] = 0.0
    g["_VENTA_MONTO_COL"] = ""
    g["_VENTA_KG_COL"] = ""

    if sales_df is None or sales_df.empty or "Cliente_ID" not in sales_df.columns:
        return g

    amount_col = detect_sales_amount_column(sales_df)
    weight_col = detect_sales_weight_column(sales_df)
    if not amount_col:
        return g

    s = sales_df.dropna(subset=["Cliente_ID"]).copy()
    if s.empty:
        return g

    if s[amount_col].dtype == object:
        s["_VENTA_MONTO"] = parse_decimal_series(s[amount_col]).fillna(0.0)
    else:
        s["_VENTA_MONTO"] = pd.to_numeric(s[amount_col], errors="coerce").fillna(0.0)

    if weight_col:
        if s[weight_col].dtype == object:
            s["_VENTA_KG"] = parse_decimal_series(s[weight_col]).fillna(0.0)
        else:
            s["_VENTA_KG"] = pd.to_numeric(s[weight_col], errors="coerce").fillna(0.0)
    else:
        s["_VENTA_KG"] = 0.0

    agg = s.groupby("Cliente_ID")[["_VENTA_MONTO", "_VENTA_KG"]].sum().reset_index()
    agg = agg.rename(columns={"_VENTA_MONTO": "FACTURACION_CLIENTE", "_VENTA_KG": "KILOS_CLIENTE"})

    out = g.merge(agg, on="Cliente_ID", how="left", suffixes=("", "_ventas"))
    if "FACTURACION_CLIENTE_ventas" in out.columns:
        out["FACTURACION_CLIENTE"] = pd.to_numeric(out["FACTURACION_CLIENTE_ventas"], errors="coerce")
        out = out.drop(columns=["FACTURACION_CLIENTE_ventas"])
    if "KILOS_CLIENTE_ventas" in out.columns:
        out["KILOS_CLIENTE"] = pd.to_numeric(out["KILOS_CLIENTE_ventas"], errors="coerce")
        out = out.drop(columns=["KILOS_CLIENTE_ventas"])
    out["FACTURACION_CLIENTE"] = pd.to_numeric(out["FACTURACION_CLIENTE"], errors="coerce").fillna(0.0)
    out["KILOS_CLIENTE"] = pd.to_numeric(out["KILOS_CLIENTE"], errors="coerce").fillna(0.0)
    out["_VENTA_MONTO_COL"] = amount_col
    out["_VENTA_KG_COL"] = weight_col or ""
    return out

def matching_mode_label(g: pd.DataFrame) -> str:
    if "_TP_MATCH_MODE" not in g.columns or g.empty:
        return ""
    vals = [norm_text(x) for x in g["_TP_MATCH_MODE"].dropna().tolist() if norm_text(x)]
    if not vals:
        return ""
    top = pd.Series(vals).mode()
    return norm_text(top.iloc[0]) if not top.empty else vals[0]

def top_missing_items_exact(g: pd.DataFrame, top_n: int = 3):
    if matching_mode_label(g) != "EXACTO":
        return []
    base = g[g["TP_ELIGIBLE"] == True].copy() if "TP_ELIGIBLE" in g.columns else g.copy()
    client_counter = Counter()
    for _, r in base.iterrows():
        skus = r.get("SKUs_faltan_para_100_ENT", [])
        if not isinstance(skus, list):
            skus = parse_listish(skus)
        seen = set()
        for sku in skus:
            key = norm_text(sku)
            if not key or key in seen:
                continue
            client_counter[key] += 1
            seen.add(key)
    rows = []
    for sku, clientes in client_counter.most_common(top_n):
        rows.append((sku, int(clientes)))
    return rows


def top_focus_items_for_day(g: pd.DataFrame, dia_code: str, top_n: int = 3):
    if matching_mode_label(g) != "EXACTO":
        return []
    base = g.copy()
    base = base[pd.to_numeric(base.get("PORTAFOLIO_PCT", np.nan), errors="coerce").notna()].copy()
    base = base[(base["PORTAFOLIO_PCT"] >= 60) & (base["PORTAFOLIO_PCT"] < 80)].copy()
    if "frecuencia" in base.columns:
        dayf = base[base["frecuencia"].apply(lambda x: freq_has_day(x, dia_code))].copy()
        if not dayf.empty:
            base = dayf
    base = base[base["Vendedor_ID"].astype(str) != "1"].copy()
    if base.empty:
        return []
    client_counter = Counter()
    for _, r in base.iterrows():
        skus = r.get("SKUs_faltan_para_80_ENT", [])
        if not isinstance(skus, list):
            skus = parse_listish(skus)
        seen = set()
        for sku in skus:
            key = norm_text(sku)
            if not key or key in seen:
                continue
            client_counter[key] += 1
            seen.add(key)
    rows = []
    for sku, clientes in client_counter.most_common(top_n):
        rows.append((sku, int(clientes)))
    return rows




VENDOR_ID_ALIASES = {
    "2": "LUDUEÑA JULIANA",
    "3": "GAMBINO NADIA",
    "4": "GRIBAUDO ANGEL",
    "5": "VERGARA MARIA JOSE",
    "6": "PEYRONEL ANDREA",
}

VENDOR_SHORT_ALIASES = {
    "2": "Juliana",
    "3": "Nadia",
    "4": "Ángel",
    "5": "Majo",
    "6": "Andrea",
}

def is_generic_vendor_label(x):
    k = norm_key(x)
    return k in {
        "",
        "SUPERVISOR",
        "PROMOTOR",
        "NO HAY VALIDACION",
        "NO HAY VALIDACI N",
        "GERENCIA",
        "GENERAL",
        "SIN VENDEDOR",
    }


def vendor_alias_key(vendedor_id):
    return norm_key(VENDOR_ID_ALIASES.get(norm_vendor_id(vendedor_id), ""))

def vendor_display_name(vendedor_id, fallback=""):
    vid = norm_vendor_id(vendedor_id)
    fb = norm_text(fallback)
    if fb and not is_generic_vendor_label(fb):
        return fb
    return VENDOR_SHORT_ALIASES.get(vid) or VENDOR_ID_ALIASES.get(vid) or (fb if fb and not is_generic_vendor_label(fb) else "") or f"Vendedor {vid}"
def norm_key(x):
    if x is None or (isinstance(x, float) and np.isnan(x)):
        return ""
    s = str(x).strip()
    s = unicodedata.normalize("NFKD", s)
    s = "".join(ch for ch in s if not unicodedata.combining(ch))
    s = s.upper()
    s = re.sub(r"[^A-Z0-9 ]+", " ", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s


def fmt_pct(n):
    try:
        return f"{float(n):.1f}%"
    except Exception:
        return "0.0%"


def fmt_int(n):
    try:
        return str(int(round(float(n))))
    except Exception:
        return "0"


def draw_metric_box_row(c, x, y_top, total_w, title, items, fill_color="#F5F8FC", title_color="#1F4E79", body_font=9.5):
    """Draw a clean one-row summary box and return consumed height in points."""
    box_h = 1.95 * cm
    c.setFillColor(colors.HexColor(fill_color))
    c.roundRect(x, y_top - box_h, total_w, box_h, 0.18 * cm, fill=1, stroke=0)
    c.setFillColor(colors.HexColor(title_color))
    c.setFont("Helvetica-Bold", 10.5)
    c.drawString(x + 0.28 * cm, y_top - 0.42 * cm, title)

    n = max(len(items), 1)
    cell_w = total_w / n
    label_y = y_top - 1.02 * cm
    value_y = y_top - 1.47 * cm
    for i, (lab, val) in enumerate(items):
        cx = x + i * cell_w
        pad = 0.18 * cm
        if i > 0:
            c.setStrokeColor(colors.HexColor("#D6E0EC"))
            c.setLineWidth(0.6)
            c.line(cx, y_top - 1.65 * cm, cx, y_top - 0.58 * cm)
        c.setFillColor(colors.HexColor("#4B5B6B"))
        c.setFont("Helvetica-Bold", 8.2)
        c.drawString(cx + pad, label_y, lab)
        c.setFillColor(colors.black)
        c.setFont("Helvetica-Bold", body_font)
        c.drawString(cx + pad, value_y, val)
    c.setStrokeColor(colors.black)
    c.setLineWidth(1)
    return box_h + 0.18 * cm


def choose_vendor_name_column(df: pd.DataFrame):
    candidates = [
        "Vendedor", "Vendedor_Nombre", "NombreVendedor", "VendedorNombre",
        "Preventista", "Preventista_Nombre", "NombrePreventista",
        "DescripcionVendedor", "DescVendedor", "Nombre_Vendedor",
        "Nombre Vendedor", "Usuario", "NombreUsuario",
    ]
    for c in candidates:
        if c in df.columns:
            return c
    for c in df.columns:
        ck = norm_key(c)
        if "VENDEDOR" in ck and "ID" not in ck:
            return c
    for c in df.columns:
        ck = norm_key(c)
        if "PREVENT" in ck and "ID" not in ck:
            return c
    return None


def build_vendor_dimension(g: pd.DataFrame) -> pd.DataFrame:
    base = g.copy()
    if "Vendedor_ID" not in base.columns:
        base["Vendedor_ID"] = np.nan

    name_col = choose_vendor_name_column(base)
    if name_col:
        base["_VEND_NOMBRE"] = base[name_col].map(norm_text)
    else:
        base["_VEND_NOMBRE"] = ""

    id_series = base["Vendedor_ID"].map(norm_vendor_id)
    base["_VEND_ID_TXT"] = id_series

    def _best_label(row):
        nm = norm_text(row.get("_VEND_NOMBRE", ""))
        vid = norm_text(row.get("_VEND_ID_TXT", ""))
        if nm and nm.lower() != "nan" and not is_generic_vendor_label(nm):
            return nm
        return vendor_display_name(vid, nm)

    base["_VEND_LABEL"] = base.apply(_best_label, axis=1)
    dim = (
        base[["_VEND_ID_TXT", "_VEND_NOMBRE", "_VEND_LABEL"]]
        .drop_duplicates()
        .copy()
    )
    dim["_MATCH_KEY"] = dim["_VEND_LABEL"].map(norm_key)
    dim["_MATCH_KEY_NAME"] = dim["_VEND_NOMBRE"].map(norm_key)
    return dim


def resolve_objective_vendor_matches(obj_vendedores: pd.DataFrame, vendor_dim: pd.DataFrame) -> pd.DataFrame:
    if obj_vendedores is None or obj_vendedores.empty:
        return pd.DataFrame(columns=["Vendedores", "OBJETIVO", "_MATCH_STATUS", "_MATCH_KEY", "_VEND_ID_TXT", "_VEND_LABEL"])

    obj = obj_vendedores.copy()
    obj["Vendedores"] = obj["Vendedores"].map(norm_text)
    obj["_MATCH_KEY"] = obj["Vendedores"].map(norm_key)
    obj["_MATCH_STATUS"] = "SIN_MATCH"
    obj["_VEND_ID_TXT"] = ""
    obj["_VEND_LABEL"] = ""

    if vendor_dim is None or vendor_dim.empty:
        return obj

    by_key = {}
    by_name = {}
    for _, r in vendor_dim.iterrows():
        if r["_MATCH_KEY"] and r["_MATCH_KEY"] not in by_key:
            by_key[r["_MATCH_KEY"]] = (r["_VEND_ID_TXT"], r["_VEND_LABEL"])
        if r["_MATCH_KEY_NAME"] and r["_MATCH_KEY_NAME"] not in by_name:
            by_name[r["_MATCH_KEY_NAME"]] = (r["_VEND_ID_TXT"], r["_VEND_LABEL"])

    all_keys = list({k for k in list(by_key.keys()) + list(by_name.keys()) if k})

    def _set_match(ix, vend_id, vend_label, status):
        obj.at[ix, "_VEND_ID_TXT"] = norm_text(vend_id)
        obj.at[ix, "_VEND_LABEL"] = norm_text(vend_label)
        obj.at[ix, "_MATCH_STATUS"] = status

    for ix, r in obj.iterrows():
        mk = r["_MATCH_KEY"]
        if not mk or mk == "TOTAL PYP":
            continue

        direct_id = ""
        m = re.fullmatch(r"VENDEDOR\s+(\d+)", mk)
        if m:
            direct_id = norm_vendor_id(m.group(1))
        elif re.fullmatch(r"\d+", mk):
            direct_id = norm_vendor_id(mk)
        if direct_id:
            _set_match(ix, direct_id, vendor_display_name(direct_id, r.get("Vendedores", "")), "ID_DIRECTO")
            continue

        if mk in by_key:
            _set_match(ix, by_key[mk][0], by_key[mk][1], "EXACTO")
            continue
        if mk in by_name:
            _set_match(ix, by_name[mk][0], by_name[mk][1], "EXACTO_NOMBRE")
            continue
        cand = difflib.get_close_matches(mk, all_keys, n=1, cutoff=0.72)
        if cand:
            hit = cand[0]
            data = by_key.get(hit) or by_name.get(hit)
            if data:
                _set_match(ix, data[0], data[1], "APROX")
    return obj


def load_objectives_tp(path_xlsx, g: pd.DataFrame, prev_day_snap=None, prev_run_snap=None):
    result = {
        "path": path_xlsx,
        "tax_df": pd.DataFrame(columns=["TAXONOMIA", "OBJETIVO", "ACUMULADO", "REAL_DIA", "FALTANTE", "CUMPLIMIENTO_PCT"]),
        "vend_df": pd.DataFrame(columns=["Vendedores", "OBJETIVO", "ACUMULADO", "REAL_DIA", "FALTANTE", "CUMPLIMIENTO_PCT", "_VEND_ID_TXT", "_VEND_LABEL", "_MATCH_STATUS"]),
        "summary_total": None,
        "loaded": False,
        "error": "",
    }
    if not path_xlsx or not os.path.exists(path_xlsx):
        result["error"] = "Archivo de objetivos no encontrado."
        return result

    try:
        raw = pd.read_excel(path_xlsx, sheet_name=0, header=None)
    except Exception as e:
        result["error"] = f"No se pudo leer Excel objetivos: {e}"
        return result

    # Bloque taxonomías: columnas B:E en la captura (índices 1:5)
    tax = raw.iloc[:, 1:5].copy()
    tax.columns = ["TAXONOMIA", "OBJETIVO", "ACUMULADO_XLSX", "FALTANTE_XLSX"]
    tax = tax.iloc[4:].copy()
    tax["TAXONOMIA"] = tax["TAXONOMIA"].map(norm_text)
    tax["OBJETIVO"] = pd.to_numeric(tax["OBJETIVO"], errors="coerce")
    tax = tax[tax["TAXONOMIA"] != ""]
    tax = tax[tax["OBJETIVO"].notna()].copy()

    # Bloque vendedores: columnas G:K (índices 6:11)
    vend = raw.iloc[:, 6:11].copy()
    vend.columns = ["Vendedores", "OBJETIVO", "ACUMULADO_XLSX", "FALTANTE_XLSX", "REAL_DIA_XLSX"]
    vend = vend.iloc[4:].copy()
    vend["Vendedores"] = vend["Vendedores"].map(norm_text)
    vend["OBJETIVO"] = pd.to_numeric(vend["OBJETIVO"], errors="coerce")
    vend = vend[vend["Vendedores"] != ""]
    vend = vend[vend["OBJETIVO"].notna()].copy()

    base = g.dropna(subset=["Cliente_ID"]).copy()
    tp_sys = base[base["TP_SISTEMA"]].copy()

    baseline_snap = prev_run_snap if prev_run_snap else prev_day_snap
    baseline_mode = "PREV_RUN" if prev_run_snap else ("PREV_DAY" if prev_day_snap else "NONE")

    # taxonomías automáticas
    cur_by_tax = tp_sys.groupby("Segmento")["Cliente_ID"].nunique().to_dict() if "Segmento" in tp_sys.columns else {}
    prev_by_tax = ((baseline_snap or {}).get("tp_sys_by_segmento") or {})
    tax["ACUMULADO"] = tax["TAXONOMIA"].map(lambda x: int(cur_by_tax.get(x, 0)))
    tax["REAL_DIA"] = tax["TAXONOMIA"].map(lambda x: int(cur_by_tax.get(x, 0)) - int(prev_by_tax.get(x, 0)))
    tax["FALTANTE"] = (tax["OBJETIVO"] - tax["ACUMULADO"]).clip(lower=0)
    tax["CUMPLIMIENTO_PCT"] = np.where(tax["OBJETIVO"] > 0, (tax["ACUMULADO"] / tax["OBJETIVO"]) * 100.0, 0.0)

    # vendedores automáticos con matching
    vendor_dim = build_vendor_dimension(base)
    vend = resolve_objective_vendor_matches(vend, vendor_dim)
    cur_by_vend = tp_sys.groupby("Vendedor_ID")["Cliente_ID"].nunique().to_dict() if "Vendedor_ID" in tp_sys.columns else {}
    cur_by_vend = {norm_vendor_id(k): int(v) for k, v in cur_by_vend.items()}
    prev_by_vend = ((baseline_snap or {}).get("tp_sys_by_vendedor") or {})
    prev_by_vend = {norm_vendor_id(k): int(v) for k, v in prev_by_vend.items()}

    def _vend_acc(row):
        k = norm_vendor_id(row.get("_VEND_ID_TXT", ""))
        return int(cur_by_vend.get(k, 0)) if k else 0

    def _vend_day(row):
        k = norm_vendor_id(row.get("_VEND_ID_TXT", ""))
        if not k:
            return 0
        return int(cur_by_vend.get(k, 0)) - int(prev_by_vend.get(k, 0))

    vend["ACUMULADO"] = vend.apply(_vend_acc, axis=1)
    vend["REAL_DIA"] = vend.apply(_vend_day, axis=1)
    vend["FALTANTE"] = (vend["OBJETIVO"] - vend["ACUMULADO"]).clip(lower=0)
    vend["CUMPLIMIENTO_PCT"] = np.where(vend["OBJETIVO"] > 0, (vend["ACUMULADO"] / vend["OBJETIVO"]) * 100.0, 0.0)

    total_obj = float(vend.loc[vend["Vendedores"].map(norm_key) == "TOTAL PYP", "OBJETIVO"].sum())
    total_acc = int(tp_sys["Cliente_ID"].nunique())
    prev_total = int((baseline_snap or {}).get("tp_sys_total", 0))
    result["summary_total"] = {
        "objetivo": total_obj if total_obj > 0 else float(tax["OBJETIVO"].sum()),
        "acumulado": total_acc,
        "real_dia": total_acc - prev_total,
        "faltante": max((total_obj if total_obj > 0 else float(tax["OBJETIVO"].sum())) - total_acc, 0),
        "cumplimiento_pct": ((total_acc / total_obj) * 100.0) if total_obj > 0 else 0.0,
        "_baseline_mode": baseline_mode,
    }
    result["tax_df"] = tax.reset_index(drop=True)
    result["vend_df"] = vend.reset_index(drop=True)
    result["loaded"] = True
    return result


def get_objetivo_vendedor_row(obj_bundle, vendedor_id, vendedor_label=""):
    if not obj_bundle or not obj_bundle.get("loaded"):
        return None
    df = obj_bundle.get("vend_df")
    if df is None or df.empty:
        return None
    vid = norm_vendor_id(vendedor_id)
    if vid:
        m = df[df["_VEND_ID_TXT"].map(norm_text) == vid]
        if not m.empty:
            return m.iloc[0].to_dict()
    vk = norm_key(vendedor_label)
    if vk:
        m = df[df["_MATCH_KEY"].map(norm_key) == vk]
        if not m.empty:
            return m.iloc[0].to_dict()
    alias_k = vendor_alias_key(vendedor_id)
    if alias_k:
        m = df[df["_MATCH_KEY"].map(norm_key) == alias_k]
        if not m.empty:
            return m.iloc[0].to_dict()
    return None

def normalize_cliente_id(series: pd.Series) -> pd.Series:
    s = series.astype(str)
    s = s.str.replace(r"\.0$", "", regex=True)
    s = s.str.extract(r"(\d+)", expand=False)
    return pd.to_numeric(s, errors="coerce").astype("Int64")


# ============================================================
# Loaders
# ============================================================

def load_sales_clean(path_csv):
    ext = os.path.splitext(path_csv)[1].lower()
    if ext in (".xlsx", ".xls"):
        df = pd.read_excel(path_csv)
    else:
        last_err = None
        df = None
        for enc in ("utf-8", "utf-8-sig", "latin1", "cp1252"):
            try:
                df = pd.read_csv(path_csv, encoding=enc, sep=None, engine="python")
                break
            except Exception as e:
                last_err = e
                df = None
        if df is None:
            raise last_err

    df.columns = [norm_text(c) for c in df.columns]

    ren = {}
    for c in df.columns:
        cl = c.lower()
        if cl == "cliente":
            ren[c] = "Cliente"
        elif cl == "localidad":
            ren[c] = "Localidad"
        elif cl in ("fecha", "fechacomprobante"):
            ren[c] = "FechaComprobante"
        elif cl == "fechaentrega":
            ren[c] = "FechaEntrega"
    df = df.rename(columns=ren)

    if "Cliente" not in df.columns:
        df["Cliente"] = np.nan
    df["Cliente_ID"] = normalize_cliente_id(df["Cliente"])

    for col in ("FechaEntrega", "FechaComprobante"):
        if col in df.columns:
            df[col] = pd.to_datetime(df[col], errors="coerce", dayfirst=True)

    if "Localidad" not in df.columns:
        df["Localidad"] = np.nan
    df["Localidad"] = df["Localidad"].astype(str).map(norm_text)

    # normalizar columnas de importe y kilos frecuentes del archivo diario
    for col in ("ImporteNetoItem", "ImporteItem", "NetoItem", "PesoKg", "Kilos", "Kg"):
        if col in df.columns:
            df[col] = parse_decimal_series(df[col])

    return df


def load_gescom(path_xlsx):
    g = pd.read_excel(path_xlsx, sheet_name="Principal")
    g.columns = [norm_text(c) for c in g.columns]

    if "Cliente" not in g.columns:
        g["Cliente"] = np.nan
    g["Cliente_ID"] = normalize_cliente_id(g["Cliente"])

    if "Razon_Social" not in g.columns and "Razon Social" in g.columns:
        g = g.rename(columns={"Razon Social": "Razon_Social"})

    for c in [
        "Vendedor_ID", "Segmento", "Razon_Social", "Direccion",
        "porcentajePortafolio", "Puntuacion",
        "SubCanalTiendaPerfecta", "TiendaPerfecta", "frecuencia"
    ]:
        if c not in g.columns:
            g[c] = np.nan

    g["porcentajePortafolio"] = to_pct_0_100(g["porcentajePortafolio"])
    g["Puntuacion"] = to_num(g["Puntuacion"])
    g["SubCanalTiendaPerfecta"] = g["SubCanalTiendaPerfecta"].astype(str).map(norm_text)
    g["TiendaPerfecta"] = g["TiendaPerfecta"].astype(str).map(norm_text)
    g["Razon_Social"] = g["Razon_Social"].astype(str).map(norm_text)
    return g

def attach_localidad_from_sales(gescom_df, sales_df):
    g = gescom_df.copy()
    if "Localidad" not in g.columns:
        g["Localidad"] = np.nan

    s = sales_df.dropna(subset=["Cliente_ID"]).copy()
    if "Localidad" not in s.columns:
        g["Localidad"] = g["Localidad"].astype(str).map(norm_text).replace("", "-").fillna("-")
        return g

    s = s[s["Localidad"].notna() & (s["Localidad"] != "")].copy()
    if s.empty:
        g["Localidad"] = g["Localidad"].astype(str).map(norm_text).replace("", "-").fillna("-")
        return g

    loc_last = s.groupby("Cliente_ID")["Localidad"].last().reset_index()
    out = g.merge(loc_last, on="Cliente_ID", how="left", suffixes=("", "_from_sales"))
    out["Localidad"] = out["Localidad_from_sales"].where(
        out["Localidad_from_sales"].notna() & (out["Localidad_from_sales"] != ""),
        out["Localidad"]
    )
    out = out.drop(columns=["Localidad_from_sales"])
    out["Localidad"] = out["Localidad"].astype(str).map(norm_text).replace("", "-").fillna("-")
    return out


# ============================================================
# TP Flags
# ============================================================
def compute_tp_flags(g):
    out = g.copy()
    out["PORTAFOLIO_PCT"] = to_pct_0_100(out["porcentajePortafolio"])
    out["TP_ELIGIBLE"] = out["SubCanalTiendaPerfecta"].astype(str).map(lambda x: norm_text(x).upper() == "SI")
    out["TP_POR_REGLA"] = out["TP_ELIGIBLE"] & (pd.to_numeric(out["Puntuacion"], errors="coerce") >= 80)
    out["TP_SISTEMA"] = out["TiendaPerfecta"].astype(str).map(lambda x: norm_text(x).upper() == "SI")
    return out


# ============================================================
# Lista de precios (XLSX/CSV)
# ============================================================
def load_price_list_from_table(path: str):
    if not path or not os.path.exists(path):
        return {"by_desc": {}, "meta": {"source": "SinListaPreciosTabla", "items": 0, "path": path}}

    ext = os.path.splitext(path)[1].lower()
    if ext in (".xlsx", ".xls"):
        df = pd.read_excel(path)
    else:
        df = pd.read_csv(path, encoding="utf-8", sep=None, engine="python")

    df.columns = [norm_text(c) for c in df.columns]

    rename = {}
    for c in df.columns:
        cl = c.lower()
        if cl in ("descripcion", "descripción", "articulo", "artículo", "producto", "item"):
            rename[c] = "Descripcion"
        if cl in ("preciounitariosiniva", "precio_unitario_sin_iva", "precio sin iva", "preciosiniva",
                  "precio unitario", "precio_unitario", "neto", "precio"):
            rename[c] = "Precio"
    df = df.rename(columns=rename)

    if "Descripcion" not in df.columns or "Precio" not in df.columns:
        return {"by_desc": {}, "meta": {"source": "ListaPreciosTabla(ColumnasNoDetectadas)", "items": 0, "path": path}}

    df["Descripcion"] = df["Descripcion"].astype(str).map(norm_text)
    df["Precio"] = pd.to_numeric(df["Precio"], errors="coerce")
    df = df[df["Descripcion"].notna() & (df["Descripcion"] != "") & df["Precio"].notna() & (df["Precio"] > 0)].copy()
    if df.empty:
        return {"by_desc": {}, "meta": {"source": "ListaPreciosTabla(Vacia)", "items": 0, "path": path}}

    grp = df.groupby("Descripcion")["Precio"].median().reset_index()
    by_desc = {norm_text(r["Descripcion"]).upper(): float(r["Precio"]) for _, r in grp.iterrows()}
    return {"by_desc": by_desc, "meta": {"source": f"ListaPreciosTabla({os.path.basename(path)})", "items": len(by_desc), "path": path}}


# ============================================================
# Universo Portafolio + Compra (EXACTO vs FALLBACK) + DEBUG
# ============================================================
STOP_TOKENS = {
    "DE","DEL","LA","LAS","EL","LOS","Y","CON","SIN","PARA","POR","X",
    "PACK","CAJA","UN","UNA","ML","GR","G","KG","LT","LTS","UNID","UNIDADES"
}
BRAND_CANON = {"DORITOS","LAY","CHEETOS","DINAMITA","PEP","PEPITAS","3D","TWISTOS","TOSTITOS","QUAKER"}

def _canon_token(t):
    if t in ("S",):
        return ""
    if t in ("LAYS","LAY","LAY'S","LAYS'","LAY’S"):
        return "LAY"
    if t in ("CHEETO","CHEETOS"):
        return "CHEETOS"
    return t

def _tokset(s):
    s = norm_text(s).upper().replace("’", "'")
    toks = re.findall(r"[A-Z0-9]+", s)
    out = set()
    for t in toks:
        if t in STOP_TOKENS or len(t) < 2:
            continue
        ct = _canon_token(t)
        if not ct or ct in STOP_TOKENS or len(ct) < 2:
            continue
        out.add(ct)
    return out

def normalize_desc(s):
    s = norm_text(s).upper().replace("’","'")
    s = re.sub(r"\s+"," ", s).strip()
    s = s.replace("LAYS","LAY").replace("LAY'S","LAY")
    return s

def extract_portafolio_universe_from_pptx(path_pptx):
    if not path_pptx or not os.path.exists(path_pptx):
        return []
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        return []
    prs = Presentation(path_pptx)
    brands = ["DORITOS","LAY","LAY'S","LAYS","CHEETOS","DINAMITA","PEP","PEPITAS","3D","TWISTOS","TOSTITOS","QUAKER"]

    cand = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            txt = shape.text_frame.text or ""
            for line in txt.splitlines():
                line = re.sub(r"\s+"," ", line.strip())
                if len(line) < 4:
                    continue
                up = line.upper()
                if any(b in up for b in brands):
                    cand.append(line)

    uni, seen = [], set()
    for it in cand:
        k = it.upper()
        if k in seen:
            continue
        seen.add(k)
        t = _tokset(it)
        if t and (t & BRAND_CANON):
            uni.append(it)
    return uni

def _score_product_column(series):
    s = series.dropna()
    if s.empty:
        return -1e9
    s = s.astype(str)
    nn = len(s)
    uniq = s.nunique()
    avg_len = float(s.map(len).mean())
    score = 0.0
    score += min(nn / 5000.0, 1.0) * 2.0
    score += min(avg_len / 25.0, 1.0) * 3.0
    score += min(uniq / max(nn, 1), 1.0) * 3.0
    num_like = s.str.fullmatch(r"\d+").mean()
    if pd.notna(num_like) and num_like > 0.6:
        score -= 3.0
    return score

def detect_best_product_column(sales):
    candidates_exact = {"articulo","artículo","producto","descripcion","descripción","item"}
    for c in sales.columns:
        if c.lower().strip() in candidates_exact:
            return c
    best_c, best_s = None, -1e18
    for c in sales.columns:
        cl = c.lower()
        if cl in ("cliente","cliente_id","localidad","fecha","fechacomprobante","fechaentrega"):
            continue
        if sales[c].dtype == "O" or "string" in str(sales[c].dtype):
            sc = _score_product_column(sales[c])
            if sc > best_s:
                best_s, best_c = sc, c
    return best_c

def build_client_purchase_tokens(sales):
    prod_col = detect_best_product_column(sales)
    if "Cliente_ID" not in sales.columns or prod_col is None:
        return {}, prod_col, []
    s = sales.dropna(subset=["Cliente_ID"]).copy()
    out = {}
    samples = []
    for cid, sub in s.groupby("Cliente_ID"):
        toks = set()
        vals = sub[prod_col].dropna().astype(str).unique().tolist()
        for v in vals:
            toks |= _tokset(normalize_desc(v))
        out[int(cid)] = toks
        if len(samples) < 5 and vals:
            samples.append(vals[0])
    return out, prod_col, samples


def _normalize_product_signature(s):
    s = norm_text(s).upper().replace("’", "'").replace("`", "'")
    s = unicodedata.normalize("NFKD", s)
    s = "".join(ch for ch in s if not unicodedata.combining(ch))
    s = re.sub(r"[^A-Z0-9' ]+", " ", s)
    s = s.replace("LAY'S", "LAYS").replace("LAY S", "LAYS")
    s = s.replace("LAYS CLASICA ", "LAYS CLASICAS ")
    s = s.replace("EXTRA FLAMIN HOT", "EXTRA FH").replace("FLAMIN HOT", "FH")
    s = s.replace("DORITOS DINAMITA", "DINAMITA")
    s = s.replace("PEPITAS", "PEP")
    s = re.sub(r"(\d{2,3})\s*GRX\d+(?:X\d+)?", r"\1G", s)
    s = re.sub(r"(\d{2,3})\s*GX\d+(?:X\d+)?", r"\1G", s)
    s = re.sub(r"(\d{2,3})\s*G\s*X\d+(?:X\d+)?", r"\1G", s)
    s = re.sub(r"(\d{2,3})X\d+(?:X\d+)?", r"\1G", s)
    s = re.sub(r"(\d{2,3})\s*GR\b", r"\1G", s)
    s = re.sub(r"\b(MB|MND|RM|PI|TIR|OT|DTS|SB)\b", " ", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s


def _product_size_token(sig):
    m = re.search(r"\b(\d{2,3})G\b", sig)
    return f"{m.group(1)}G" if m else ""


def _product_tokens(sig):
    toks = []
    for t in re.findall(r"[A-Z0-9]+", sig):
        if t in STOP_TOKENS or t in {"MB", "MND", "RM", "PI", "TIR", "OT", "DTS", "SB"}:
            continue
        if re.fullmatch(r"\d+", t):
            continue
        if re.fullmatch(r"X\d+", t):
            continue
        ct = _canon_token(t)
        if not ct:
            continue
        if len(ct) < 2 and ct != "3D":
            continue
        toks.append(ct)
    return toks


def _product_brand(tokens):
    for t in tokens:
        if t in BRAND_CANON:
            return t
    return ""


def _product_match_metrics(port_item, sold_item):
    p_sig = _normalize_product_signature(port_item)
    s_sig = _normalize_product_signature(sold_item)
    p_tokens = set(_product_tokens(p_sig))
    s_tokens = set(_product_tokens(s_sig))
    p_brand = _product_brand(list(p_tokens))
    s_brand = _product_brand(list(s_tokens))
    p_size = _product_size_token(p_sig)
    s_size = _product_size_token(s_sig)

    brand_ok = (not p_brand) or (not s_brand) or (p_brand == s_brand)
    size_ok = (not p_size) or (not s_size) or (p_size == s_size)

    p_core = {t for t in p_tokens if t not in {p_brand, p_size}}
    s_core = {t for t in s_tokens if t not in {s_brand, s_size}}

    shared_core = len(p_core & s_core)
    ratio = difflib.SequenceMatcher(None, p_sig, s_sig).ratio()

    score = ratio * 100.0
    if brand_ok and p_brand and s_brand:
        score += 16.0
    if size_ok and p_size and s_size:
        score += 14.0
    score += shared_core * 24.0

    return {
        "score": score,
        "ratio": ratio,
        "shared_core": shared_core,
        "brand_ok": brand_ok,
        "size_ok": size_ok,
        "p_sig": p_sig,
        "s_sig": s_sig,
    }


def _product_is_match(port_item, sold_item):
    m = _product_match_metrics(port_item, sold_item)
    if not m["brand_ok"] or not m["size_ok"]:
        return False
    if m["shared_core"] >= 2:
        return True
    if m["shared_core"] >= 1 and m["ratio"] >= 0.58:
        return True
    if m["ratio"] >= 0.90:
        return True
    return False


def _build_client_product_records(sales):
    prod_col = detect_best_product_column(sales)
    if "Cliente_ID" not in sales.columns or prod_col is None:
        return {}, prod_col, [], []

    s = sales.dropna(subset=["Cliente_ID"]).copy()
    by_client = {}
    samples = []
    catalog = []

    for cid, sub in s.groupby("Cliente_ID"):
        vals = [norm_text(v) for v in sub[prod_col].dropna().astype(str).unique().tolist() if norm_text(v)]
        by_client[int(cid)] = vals
        if len(samples) < 5:
            samples.extend(vals[: max(0, 5 - len(samples))])
        catalog.extend(vals)

    catalog = list(dict.fromkeys(catalog))
    return by_client, prod_col, samples[:5], catalog


def attach_portafolio_missing_to_g(g, sales, pptx_path, debug_path):
    out = g.copy()
    universe = extract_portafolio_universe_from_pptx(pptx_path)
    universe = [norm_text(u) for u in universe if norm_text(u)]
    universe = list(dict.fromkeys(universe))

    u_size = len(universe)
    need80 = int(ceil(0.80 * u_size)) if u_size else 0

    client_products, prod_col, samples, catalog = _build_client_product_records(sales)
    nonempty_clients = sum(1 for _, vals in client_products.items() if vals)

    global_hits = 0
    if u_size and catalog:
        for sku in universe:
            if any(_product_is_match(sku, art) for art in catalog):
                global_hits += 1

    mode = "EXACTO" if (u_size > 0 and global_hits > 0) else ("SIN_MATCH" if u_size > 0 else "SIN_PORTAFOLIO")
    out["_TP_MATCH_MODE"] = mode
    out["_UNIVERSE_SIZE"] = u_size

    safe_mkdir(os.path.dirname(debug_path))
    with open(debug_path, "w", encoding="utf-8") as f:
        f.write(f"DEBUG TP COMPRA - {datetime.now().isoformat()}\n")
        f.write(f"portafolio_path: {pptx_path}\n")
        f.write(f"prod_col_detectado: {prod_col}\n")
        f.write(f"clientes_con_productos: {nonempty_clients}\n")
        f.write(f"universe_size: {u_size} | need80: {need80}\n")
        f.write(f"global_hits_universe_vs_compra: {global_hits}\n")
        f.write(f"modo_usado: {mode}\n")
        f.write("samples_producto:\n")
        for s in samples:
            f.write(f"  - {s}\n")

    f80_all, f100_all, n80_all, n100_all = [], [], [], []
    for _, r in out.iterrows():
        port = pd.to_numeric(r.get("PORTAFOLIO_PCT"), errors="coerce")
        cid = r.get("Cliente_ID")
        cid = int(cid) if not pd.isna(cid) else None

        if u_size == 0 or pd.isna(port):
            f80_all.append([])
            f100_all.append([])
            n80_all.append(0)
            n100_all.append(0)
            continue

        if mode == "EXACTO":
            sold_items = client_products.get(cid, [])
            missing = []
            covered = 0
            for sku in universe:
                matched = False
                for sold in sold_items:
                    if _product_is_match(sku, sold):
                        matched = True
                        break
                if matched:
                    covered += 1
                else:
                    missing.append(sku)

            faltan100 = len(missing)
            faltan80 = max(need80 - covered, 0)
            f80 = missing[:faltan80] if faltan80 else []
            f80_all.append(f80)
            f100_all.append(missing)
            n80_all.append(int(faltan80))
            n100_all.append(int(faltan100))
        else:
            covered_est = int(round((float(port) / 100.0) * u_size))
            covered_est = max(0, min(u_size, covered_est))
            faltan100 = u_size - covered_est
            faltan80 = max(need80 - covered_est, 0)
            f80_all.append(universe[:] if faltan80 else [])
            f100_all.append(universe[:])
            n80_all.append(int(faltan80))
            n100_all.append(int(faltan100))

    out["FALTAN_80_N"] = n80_all
    out["FALTAN_100_N"] = n100_all
    out["SKUs_faltan_para_80_ENT"] = f80_all
    out["SKUs_faltan_para_100_ENT"] = f100_all

    # Recalculate PORTAFOLIO_PCT from actual product coverage (replaces Gescom value)
    # Gescom may report 100% while the motor detected missing products — use motor truth.
    if u_size > 0:
        covered_n = u_size - out["FALTAN_100_N"].astype(int)
        out["PORTAFOLIO_PCT"] = (covered_n / u_size * 100).round(0).clip(0, 100).astype(int)

    return out


# ============================================================
# Historial (deltas)
# ============================================================
def load_hist_runs(hist_path):
    if not os.path.exists(hist_path):
        return []
    try:
        with open(hist_path, "r", encoding="utf-8") as f:
            obj = json.load(f)
        return obj if isinstance(obj, list) else []
    except Exception:
        return []

def save_hist_runs(hist_path, runs):
    safe_mkdir(os.path.dirname(hist_path))
    with open(hist_path, "w", encoding="utf-8") as f:
        json.dump(runs, f, ensure_ascii=False, indent=2)

def build_snapshot(g):
    base = g.dropna(subset=["Cliente_ID"]).copy()
    tp_sys = base[base["TP_SISTEMA"]].copy()
    total = int(tp_sys["Cliente_ID"].nunique())
    by_seg = tp_sys.groupby("Segmento")["Cliente_ID"].nunique().to_dict()
    by_vend = tp_sys.groupby("Vendedor_ID")["Cliente_ID"].nunique().to_dict()
    return {
        "tp_sys_total": total,
        "tp_sys_by_segmento": {str(k): int(v) for k, v in by_seg.items()},
        "tp_sys_by_vendedor": {str(k): int(v) for k, v in by_vend.items()},
    }


def diff_snap(cur, prev):
    if not cur or not prev:
        return None

    def _diff_map(a, b):
        keys = set((a or {}).keys()) | set((b or {}).keys())
        out = {}
        for k in keys:
            out[str(k)] = int((a or {}).get(k, 0)) - int((b or {}).get(k, 0))
        return out

    return {
        "tp_sys_total": int(cur.get("tp_sys_total", 0)) - int(prev.get("tp_sys_total", 0)),
        "tp_sys_by_segmento": _diff_map(cur.get("tp_sys_by_segmento"), prev.get("tp_sys_by_segmento")),
        "tp_sys_by_vendedor": _diff_map(cur.get("tp_sys_by_vendedor"), prev.get("tp_sys_by_vendedor")),
    }

def update_hist_and_get_deltas(g, hist_runs, hist_path):
    now_ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    today_str = date.today().isoformat()
    snap = build_snapshot(g)

    prev_run_snap = hist_runs[-1]["snapshot"] if (hist_runs and "snapshot" in hist_runs[-1]) else None
    delta_prev_run = diff_snap(snap, prev_run_snap) if prev_run_snap else None

    prev_day_snap = None
    for r in reversed(hist_runs):
        ts = str(r.get("ts",""))
        if ts[:10] < today_str and "snapshot" in r:
            prev_day_snap = r["snapshot"]
            break
    delta_prev_day = diff_snap(snap, prev_day_snap) if prev_day_snap else None

    hist_runs.append({"ts": now_ts, "snapshot": snap})
    if len(hist_runs) > 400:
        hist_runs = hist_runs[-400:]
    save_hist_runs(hist_path, hist_runs)
    return hist_runs, snap, delta_prev_run, delta_prev_day


# ============================================================
# Potencial (60–79) con lista precios
# ============================================================
def match_price_for_sku_text(sku_text: str, price_by_desc: dict, cache: dict):
    raw = norm_text(sku_text)
    if not raw:
        return None
    key = raw.upper()
    if key in cache:
        return cache[key]

    sku_norm = normalize_desc(raw)
    tsku = _tokset(sku_norm)

    candidates = []
    for desc_up, price in price_by_desc.items():
        dnorm = normalize_desc(desc_up)
        tdesc = _tokset(dnorm)
        shared = len(tsku & tdesc)
        if shared >= 2:
            candidates.append((desc_up, price, shared))

    if candidates:
        best = max(candidates, key=lambda x: x[2])
        cache[key] = best[1]
        return best[1]

    # fallback similitud
    best_desc, best_ratio = None, 0.0
    for desc_up in price_by_desc.keys():
        r = difflib.SequenceMatcher(None, sku_norm, normalize_desc(desc_up)).ratio()
        if r > best_ratio:
            best_ratio = r
            best_desc = desc_up
    if best_desc is not None and best_ratio >= 0.72:
        cache[key] = price_by_desc.get(best_desc)
        return cache[key]

    cache[key] = None
    return None

def compute_potential_60_79_for_tomorrow(g: pd.DataFrame, dia_code: str, price_list: dict):
    by_desc = price_list.get("by_desc", {}) or {}
    meta = price_list.get("meta", {}) or {}

    if not by_desc:
        return {"has_price_data": False, "pot_total": 0.0, "unknown": 0, "valued": 0, "falt": 0,
                "source": meta.get("source",""), "items": int(meta.get("items",0))}

    opp = g[(pd.to_numeric(g["PORTAFOLIO_PCT"], errors="coerce").notna()) &
            (g["PORTAFOLIO_PCT"] >= 60) & (g["PORTAFOLIO_PCT"] < 80)].copy()

    if "frecuencia" in opp.columns:
        od = opp[opp["frecuencia"].apply(lambda x: freq_has_day(x, dia_code))].copy()
        if not od.empty:
            opp = od

    opp = opp[opp["Vendedor_ID"].astype(str) != "1"].copy()
    if opp.empty:
        return {"has_price_data": True, "pot_total": 0.0, "unknown": 0, "valued": 0, "falt": 0,
                "source": meta.get("source",""), "items": int(meta.get("items",0))}

    cache = {}
    pot_total = 0.0
    unknown = 0
    valued = 0
    falt = 0

    for _, r in opp.iterrows():
        falt80 = r.get("SKUs_faltan_para_80_ENT", [])
        if not isinstance(falt80, list):
            falt80 = parse_listish(falt80)
        falt += len(falt80)
        for sku in falt80:
            p = match_price_for_sku_text(sku, by_desc, cache)
            if p is None:
                unknown += 1
            else:
                pot_total += float(p)
                valued += 1

    return {"has_price_data": True, "pot_total": float(pot_total), "unknown": int(unknown), "valued": int(valued),
            "falt": int(falt), "source": meta.get("source",""), "items": int(meta.get("items",0))}


# ============================================================
# PDF Vendedores (orden pedido + separador OK)
# ============================================================

def export_pdf_vendedores(g, dia_code, out_dir, obj_bundle=None):
    safe_mkdir(out_dir)

    col_cliente = "Cliente_ID"
    col_rs = "Razon_Social"
    col_port = "PORTAFOLIO_PCT"
    col_freq = "frecuencia"

    gg = g.copy()
    gg["_EN_DIA"] = True
    if col_freq in gg.columns:
        gg["_EN_DIA"] = gg[col_freq].apply(lambda x: freq_has_day(x, dia_code))

    name_col = choose_vendor_name_column(gg)
    vendedores = [v for v in gg["Vendedor_ID"].dropna().unique().tolist() if str(v) != "1"]
    created = []

    for v in vendedores:
        gv = gg[(gg["Vendedor_ID"] == v) & (gg["_EN_DIA"])].copy()
        if gv.empty:
            gv = gg[gg["Vendedor_ID"] == v].copy()
        if gv.empty:
            continue

        vend_label = ""
        if name_col and name_col in gv.columns:
            vals = [norm_text(x) for x in gv[name_col].dropna().tolist() if norm_text(x)]
            if vals:
                vend_label = vals[0]
        vend_label = vendor_display_name(v, vend_label)

        port_prom = float(pd.to_numeric(gv[col_port], errors="coerce").mean()) if pd.to_numeric(gv[col_port], errors="coerce").notna().any() else 0.0
        u_size = int(pd.to_numeric(gv.get("_UNIVERSE_SIZE", 0), errors="coerce").max()) if "_UNIVERSE_SIZE" in gv.columns else 0
        obj_row = get_objetivo_vendedor_row(obj_bundle, v, vend_label)

        pdf_path = os.path.join(out_dir, f"Vendedor_{v}_{dia_code}.pdf")
        c = canvas.Canvas(pdf_path, pagesize=landscape(A4))
        W, H = landscape(A4)
        margin = 1.4 * cm
        y = H - margin

        def new_page():
            nonlocal y
            finalize_page(c, W, H)
            c.showPage()
            y = H - margin

        # Header
        c.setFillColor(colors.HexColor("#1F4E79"))
        c.rect(margin, y - 1.1 * cm, W - 2 * margin, 1.1 * cm, fill=1, stroke=0)
        c.setFillColor(colors.white)
        c.setFont("Helvetica-Bold", 14)
        c.drawString(margin + 0.4 * cm, y - 0.75 * cm, f"Tienda Perfecta - Ruta {dia_code} - Vendedor {vend_label}")
        c.setFont("Helvetica", 9)
        c.drawRightString(W - margin - 0.4 * cm, y - 0.75 * cm, f"Generado: {date.today().isoformat()}")
        y -= 1.55 * cm

        c.setFillColor(colors.black)
        c.setFont("Helvetica", 10)
        c.drawString(margin, y, f"Clientes del día: {gv[col_cliente].nunique()} | Prom. Portafolio: {port_prom:.1f}% | Universo: {u_size}")
        y -= 0.65 * cm

        if obj_row:
            items = [
                ("Objetivo", fmt_int(obj_row.get("OBJETIVO", 0))),
                ("Acumulado", fmt_int(obj_row.get("ACUMULADO", 0))),
                ("Real día", fmt_int(obj_row.get("REAL_DIA", 0))),
                ("Faltante", fmt_int(obj_row.get("FALTANTE", 0))),
                ("Cumplimiento", fmt_pct(obj_row.get("CUMPLIMIENTO_PCT", 0))),
            ]
            consumed = draw_metric_box_row(c, margin, y, W - 2 * margin, "Resumen objetivo TP", items)
            if norm_text(obj_row.get("_MATCH_STATUS", "")) not in ("", "EXACTO", "EXACTO_NOMBRE", "ID_DIRECTO"):
                c.setFillColor(colors.HexColor("#7A5B00"))
                c.setFont("Helvetica", 7.5)
                c.drawRightString(W - margin - 0.25 * cm, y - 1.72 * cm, f"Match objetivo: {norm_text(obj_row.get('_MATCH_STATUS',''))}")
                c.setFillColor(colors.black)
            y -= consumed
        else:
            c.setFillColor(colors.HexColor("#FFF7E6"))
            c.roundRect(margin, y - 1.0 * cm, W - 2 * margin, 0.9 * cm, 0.18 * cm, fill=1, stroke=0)
            c.setFillColor(colors.HexColor("#7A5B00"))
            c.setFont("Helvetica", 8.5)
            c.drawString(margin + 0.25 * cm, y - 0.55 * cm, "Objetivo TP no encontrado para este vendedor en objetivo_vendedor_tp.xlsx.")
            c.setFillColor(colors.black)
            y -= 1.15 * cm

        c.setStrokeColor(colors.HexColor("#C9D7E6"))
        c.line(margin, y, W - margin, y)
        c.setStrokeColor(colors.black)
        y -= 0.35 * cm

        def wrap_items_to_lines(items, max_width_pts, font_name, font_size):
            if not items:
                return [""]
            lines, cur = [], ""
            for it in items:
                it = norm_text(it)
                if not it:
                    continue
                cand = it if cur == "" else f"{cur}, {it}"
                if c.stringWidth(cand, font_name, font_size) <= max_width_pts:
                    cur = cand
                else:
                    if cur:
                        lines.append(cur)
                    cur = it
            if cur:
                lines.append(cur)
            return lines

        p = pd.to_numeric(gv[col_port], errors="coerce")
        band_60_79 = gv[(p >= 60) & (p <= 79)].copy()
        band_30_59 = gv[(p >= 30) & (p < 60)].copy()
        band_lt30  = gv[p < 30].copy()

        def draw_section(title, df):
            nonlocal y
            c.setFont("Helvetica-Bold", 10)
            c.setFillColor(colors.HexColor("#1F4E79"))
            c.drawString(margin, y, f"{title} (clientes: {df[col_cliente].nunique()})")
            y -= 0.55 * cm

            if df.empty:
                c.setFillColor(colors.black)
                c.setFont("Helvetica", 8.5)
                c.drawString(margin + 0.2 * cm, y, "Sin clientes en esta banda.")
                y -= 0.65 * cm
                return

            headers = [
                ("Cliente", 1.6 * cm),
                ("Razón social", 8.6 * cm),
                ("Port.%", 1.3 * cm),
                ("Faltan80", 1.4 * cm),
                ("Faltan100", 1.6 * cm),
                ("Marcas (faltantes p/100)", 10.8 * cm),
            ]

            c.setFillColor(colors.HexColor("#E7EEF7"))
            c.rect(margin, y - 0.55 * cm, W - 2 * margin, 0.55 * cm, fill=1, stroke=0)
            c.setFillColor(colors.black)
            c.setFont("Helvetica-Bold", 9)
            x = margin
            for hname, hw in headers:
                c.drawString(x + 0.1 * cm, y - 0.4 * cm, hname)
                x += hw
            y -= 0.8 * cm

            body_font = "Helvetica"
            body_size = 7.4
            line_gap = 0.36 * cm
            min_row_h = 0.60 * cm

            df2 = df.copy()
            df2["_p"] = pd.to_numeric(df2[col_port], errors="coerce")
            df2["FALTAN_80_N"] = pd.to_numeric(df2.get("FALTAN_80_N", 0), errors="coerce").fillna(0).astype(int)
            df2["FALTAN_100_N"] = pd.to_numeric(df2.get("FALTAN_100_N", 0), errors="coerce").fillna(0).astype(int)
            df2 = df2.sort_values(by=["FALTAN_80_N", "_p"], ascending=[True, False])

            resumen_width_pts = float(headers[-1][1])

            for _, r in df2.iterrows():
                f100 = r.get("SKUs_faltan_para_100_ENT", [])
                if not isinstance(f100, list):
                    f100 = parse_listish(f100)

                n80 = int(r.get("FALTAN_80_N", 0) or 0)
                n100 = int(r.get("FALTAN_100_N", 0) or 0)

                f100_show = f100[:max(n100, 0)] if n100 > 0 else []
                lines = wrap_items_to_lines(f100_show, resumen_width_pts - 0.25 * cm, body_font, body_size)
                if len(lines) > 4:
                    lines = lines[:3] + [lines[3] + " …"]

                row_h = max(min_row_h, len(lines) * line_gap + 0.15 * cm)
                if y < margin + 2.0 * cm + row_h:
                    new_page()

                col_x = [margin]
                for _, hw in headers:
                    col_x.append(col_x[-1] + hw)

                baseline = y - 0.30 * cm
                c.setFillColor(colors.black)
                c.setFont(body_font, body_size)

                cliente_txt = "" if pd.isna(r.get(col_cliente)) else str(int(r.get(col_cliente)))
                rs_txt = norm_text(r.get(col_rs, ""))[:120]
                port_txt = "" if pd.isna(r.get("_p")) else f"{float(r.get('_p')):.1f}"

                c.drawString(col_x[0] + 0.1 * cm, baseline, cliente_txt)
                c.drawString(col_x[1] + 0.1 * cm, baseline, rs_txt)
                c.drawString(col_x[2] + 0.1 * cm, baseline, port_txt)
                c.drawString(col_x[3] + 0.1 * cm, baseline, str(n80))
                c.drawString(col_x[4] + 0.1 * cm, baseline, str(n100))

                yy = baseline
                for ln in lines:
                    c.drawString(col_x[5] + 0.1 * cm, yy, ln)
                    yy -= line_gap

                bottom = y - row_h
                y = bottom

                c.setStrokeColor(colors.lightgrey)
                c.line(margin, y + 0.08 * cm, W - margin, y + 0.08 * cm)
                c.setStrokeColor(colors.black)

            y -= 0.35 * cm

        draw_section("1) Portafolio 60% a 79% (zona de oportunidad)", band_60_79)
        draw_section("2) Portafolio 30% a 59% (prioridad de recuperación)", band_30_59)
        draw_section("3) Portafolio < 30% (crítico)", band_lt30)

        finalize_page(c, W, H)
        c.save()
        created.append(pdf_path)

    return created


# ============================================================
# PDF Gerencial COMPLETO (restaurado)
# ============================================================


def export_pdf_gerencial(g, snap, delta_prev_run, delta_prev_day, dia_code, out_dir, price_list, obj_bundle=None):
    safe_mkdir(out_dir)
    pdf_path = os.path.join(out_dir, f"RESUMEN_GERENCIAL_{date.today().isoformat()}.pdf")

    c = canvas.Canvas(pdf_path, pagesize=A4)
    W, H = A4
    margin = 1.75 * cm
    y = H - margin

    def new_page():
        nonlocal y
        finalize_page(c, W, H)
        c.showPage()
        y = H - margin

    def ensure_space(height_needed):
        nonlocal y
        if y < margin + height_needed:
            new_page()

    # Header Orbit
    c.setFillColor(colors.HexColor("#1F4E79"))
    c.roundRect(margin, y - 1.55 * cm, W - 2 * margin, 1.55 * cm, 0.18 * cm, fill=1, stroke=0)
    c.setFillColor(colors.white)
    c.setFont("Helvetica-Bold", 15)
    c.drawString(margin + 0.35 * cm, y - 0.68 * cm, "TIENDA PERFECTA ANALYTICS")
    c.setFont("Helvetica", 9.5)
    c.drawString(margin + 0.35 * cm, y - 1.10 * cm, "Powered by ORBIT")
    c.setFont("Helvetica", 8.8)
    c.drawRightString(W - margin - 0.35 * cm, y - 0.68 * cm, f"Zona: {dia_code}")
    c.drawRightString(W - margin - 0.35 * cm, y - 1.10 * cm, f"Generado: {date.today().isoformat()}")
    y -= 1.95 * cm

    if obj_bundle and obj_bundle.get("loaded"):
        summ = obj_bundle.get("summary_total") or {}
        items = [
            ("Objetivo", fmt_int(summ.get("objetivo", 0))),
            ("Acumulado", fmt_int(summ.get("acumulado", 0))),
            ("Real día", fmt_int(summ.get("real_dia", 0))),
            ("Faltante", fmt_int(summ.get("faltante", 0))),
            ("Cumplimiento", fmt_pct(summ.get("cumplimiento_pct", 0))),
        ]
        y -= draw_metric_box_row(c, margin, y, W - 2 * margin, "Resumen objetivo TP distribuidora", items)

        tax_df = obj_bundle.get("tax_df")
        if tax_df is not None and not tax_df.empty:
            ensure_space(3.3 * cm)
            c.setFont("Helvetica-Bold", 11)
            c.setFillColor(colors.black)
            c.drawString(margin, y, "1) Objetivo TP por taxonomía")
            y -= 0.55 * cm

            headers = [("Tax.", 2.0 * cm), ("Obj.", 1.8 * cm), ("Acum.", 1.8 * cm), ("Real día", 1.8 * cm), ("Faltan", 1.9 * cm), ("% Cumpl.", 2.2 * cm)]
            c.setFillColor(colors.HexColor("#E7EEF7"))
            c.rect(margin, y - 0.5 * cm, W - 2 * margin, 0.5 * cm, fill=1, stroke=0)
            c.setFillColor(colors.black)
            c.setFont("Helvetica-Bold", 9)
            x = margin
            for hname, hw in headers:
                c.drawString(x + 0.08 * cm, y - 0.34 * cm, hname)
                x += hw
            y -= 0.68 * cm

            for _, r in tax_df.iterrows():
                ensure_space(0.55 * cm)
                x = margin
                vals = [
                    norm_text(r.get("TAXONOMIA", "")),
                    fmt_int(r.get("OBJETIVO", 0)),
                    fmt_int(r.get("ACUMULADO", 0)),
                    fmt_int(r.get("REAL_DIA", 0)),
                    fmt_int(r.get("FALTANTE", 0)),
                    fmt_pct(r.get("CUMPLIMIENTO_PCT", 0)),
                ]
                widths = [2.0 * cm, 1.8 * cm, 1.8 * cm, 1.8 * cm, 1.9 * cm, 2.2 * cm]
                c.setFont("Helvetica", 9)
                for val, w in zip(vals, widths):
                    c.drawString(x + 0.08 * cm, y, val)
                    x += w
                y -= 0.42 * cm
            y -= 0.25 * cm

        vend_df = obj_bundle.get("vend_df")
        if vend_df is not None and not vend_df.empty:
            ensure_space(4.0 * cm)
            c.setFont("Helvetica-Bold", 11)
            c.drawString(margin, y, "2) Objetivo TP aperturado por vendedor")
            y -= 0.55 * cm

            headers = [("Vendedor", 6.1 * cm), ("Obj.", 1.5 * cm), ("Acum.", 1.6 * cm), ("Real día", 1.7 * cm), ("Faltan", 1.7 * cm), ("% Cumpl.", 1.9 * cm)]
            c.setFillColor(colors.HexColor("#E7EEF7"))
            c.rect(margin, y - 0.5 * cm, W - 2 * margin, 0.5 * cm, fill=1, stroke=0)
            c.setFillColor(colors.black)
            c.setFont("Helvetica-Bold", 8.5)
            x = margin
            for hname, hw in headers:
                c.drawString(x + 0.08 * cm, y - 0.34 * cm, hname)
                x += hw
            y -= 0.68 * cm

            show_df = vend_df[vend_df["Vendedores"].map(norm_key) != "TOTAL PYP"].copy()
            for _, r in show_df.iterrows():
                ensure_space(0.55 * cm)
                x = margin
                vend_name = vendor_display_name(r.get("_VEND_ID_TXT", ""), r.get("_VEND_LABEL", "") or r.get("Vendedores", ""))
                vals = [
                    norm_text(vend_name)[:42],
                    fmt_int(r.get("OBJETIVO", 0)),
                    fmt_int(r.get("ACUMULADO", 0)),
                    fmt_int(r.get("REAL_DIA", 0)),
                    fmt_int(r.get("FALTANTE", 0)),
                    fmt_pct(r.get("CUMPLIMIENTO_PCT", 0)),
                ]
                widths = [6.1 * cm, 1.5 * cm, 1.6 * cm, 1.7 * cm, 1.7 * cm, 1.9 * cm]
                c.setFont("Helvetica", 8.7)
                for val, w in zip(vals, widths):
                    c.drawString(x + 0.08 * cm, y, val)
                    x += w
                y -= 0.42 * cm
            y -= 0.30 * cm

    # 3) Portafolio distribuidora
    ensure_space(2.6 * cm)
    port_prom_total = float(pd.to_numeric(g["PORTAFOLIO_PCT"], errors="coerce").mean()) if pd.to_numeric(g["PORTAFOLIO_PCT"], errors="coerce").notna().any() else 0.0
    g_tp_aplica = g[g["TP_ELIGIBLE"] == True].copy() if "TP_ELIGIBLE" in g.columns else g.iloc[0:0].copy()
    port_prom_tp = float(pd.to_numeric(g_tp_aplica["PORTAFOLIO_PCT"], errors="coerce").mean()) if ("PORTAFOLIO_PCT" in g_tp_aplica.columns and pd.to_numeric(g_tp_aplica["PORTAFOLIO_PCT"], errors="coerce").notna().any()) else 0.0
    c.setFillColor(colors.black)
    c.setFont("Helvetica-Bold", 11)
    c.drawString(margin, y, "3) Portafolio distribuidora")
    y -= 0.6 * cm
    c.setFont("Helvetica", 10)
    c.drawString(margin + 0.35 * cm, y, f"Promedio Portafolio (todos los clientes): {port_prom_total:.1f}%")
    y -= 0.45 * cm
    c.drawString(margin + 0.35 * cm, y, f"Promedio Portafolio clientes TP (aplican): {port_prom_tp:.1f}%")
    y -= 0.85 * cm

    # 4) Facturación TP vs No TP
    ensure_space(5.2 * cm)
    c.setFont("Helvetica-Bold", 11)
    c.drawString(margin, y, "4) Facturación clientes TP vs No TP")
    y -= 0.6 * cm

    fact_col_exists = "FACTURACION_CLIENTE" in g.columns
    fact_series = pd.to_numeric(g.get("FACTURACION_CLIENTE", 0), errors="coerce").fillna(0.0) if fact_col_exists else pd.Series(dtype=float)
    kilos_series = pd.to_numeric(g.get("KILOS_CLIENTE", 0), errors="coerce").fillna(0.0) if "KILOS_CLIENTE" in g.columns else pd.Series(dtype=float)
    if fact_col_exists and float(fact_series.abs().sum()) > 0:
        g_tp_sistema = g[g["TP_SISTEMA"] == True].copy() if "TP_SISTEMA" in g.columns else g.iloc[0:0].copy()
        g_no_tp = g[g["TP_SISTEMA"] != True].copy() if "TP_SISTEMA" in g.columns else g.copy()

        fact_tp = float(pd.to_numeric(g_tp_sistema["FACTURACION_CLIENTE"], errors="coerce").fillna(0.0).sum()) if not g_tp_sistema.empty else 0.0
        fact_no_tp = float(pd.to_numeric(g_no_tp["FACTURACION_CLIENTE"], errors="coerce").fillna(0.0).sum()) if not g_no_tp.empty else 0.0
        fact_total = fact_tp + fact_no_tp

        pct_tp = (fact_tp / fact_total * 100.0) if fact_total > 0 else 0.0
        pct_no_tp = (fact_no_tp / fact_total * 100.0) if fact_total > 0 else 0.0

        clientes_tp = int(g_tp_sistema["Cliente_ID"].nunique()) if "Cliente_ID" in g_tp_sistema.columns else 0
        clientes_no_tp = int(g_no_tp["Cliente_ID"].nunique()) if "Cliente_ID" in g_no_tp.columns else 0
        prom_fact_tp = (fact_tp / clientes_tp) if clientes_tp > 0 else 0.0
        prom_fact_no_tp = (fact_no_tp / clientes_no_tp) if clientes_no_tp > 0 else 0.0

        kilos_tp = float(pd.to_numeric(g_tp_sistema.get("KILOS_CLIENTE", 0), errors="coerce").fillna(0.0).sum()) if not g_tp_sistema.empty else 0.0
        kilos_no_tp = float(pd.to_numeric(g_no_tp.get("KILOS_CLIENTE", 0), errors="coerce").fillna(0.0).sum()) if not g_no_tp.empty else 0.0
        prom_kilos_tp = (kilos_tp / clientes_tp) if clientes_tp > 0 else 0.0
        prom_kilos_no_tp = (kilos_no_tp / clientes_no_tp) if clientes_no_tp > 0 else 0.0

        delta_fact_pct = ((prom_fact_tp / prom_fact_no_tp) - 1.0) * 100.0 if prom_fact_no_tp > 0 else 0.0
        delta_kg_pct = ((prom_kilos_tp / prom_kilos_no_tp) - 1.0) * 100.0 if prom_kilos_no_tp > 0 else 0.0

        c.setFont("Helvetica", 10)
        c.drawString(margin + 0.35 * cm, y, f"Clientes TP (TiendaPerfecta=SI): {fmt_pesos(fact_tp)} | {pct_tp:.1f}%")
        y -= 0.45 * cm
        c.drawString(margin + 0.35 * cm, y, f"Clientes No TP: {fmt_pesos(fact_no_tp)} | {pct_no_tp:.1f}%")
        y -= 0.52 * cm

        c.setFont("Helvetica-Bold", 10)
        c.drawString(margin + 0.35 * cm, y, "Promedio de compra por cliente")
        y -= 0.42 * cm
        c.setFont("Helvetica", 10)
        c.drawString(margin + 0.55 * cm, y, f"Cliente TP: {fmt_pesos(prom_fact_tp)} | {prom_kilos_tp:.1f} kg")
        y -= 0.42 * cm
        c.drawString(margin + 0.55 * cm, y, f"Cliente No TP: {fmt_pesos(prom_fact_no_tp)} | {prom_kilos_no_tp:.1f} kg")
        y -= 0.42 * cm
        c.drawString(margin + 0.55 * cm, y, f"Cliente TP compra: {delta_fact_pct:+.1f}% en $ | {delta_kg_pct:+.1f}% en kg")
        y -= 0.42 * cm

        monto_col = norm_text(g["_VENTA_MONTO_COL"].iloc[0]) if "_VENTA_MONTO_COL" in g.columns and not g.empty else ""
        kg_col = norm_text(g["_VENTA_KG_COL"].iloc[0]) if "_VENTA_KG_COL" in g.columns and not g.empty else ""
        fuente_txt = monto_col or ""
        if kg_col:
            fuente_txt = f"{fuente_txt} | KG: {kg_col}" if fuente_txt else f"KG: {kg_col}"
        if fuente_txt:
            c.setFont("Helvetica-Oblique", 8.2)
            c.setFillColor(colors.HexColor("#4B5B6B"))
            c.drawString(margin + 0.35 * cm, y, f"Fuente ventas diaria: {fuente_txt}")
            c.setFillColor(colors.black)
            y -= 0.40 * cm
    else:
        c.setFont("Helvetica", 10)
        c.drawString(margin + 0.35 * cm, y, "No se pudo calcular la facturación con el archivo de ventas de esta corrida.")
        y -= 0.45 * cm
    y -= 0.20 * cm

    # 5) Oportunidad
    ensure_space(4.2 * cm)
    c.setFont("Helvetica-Bold", 11)
    c.drawString(margin, y, f"5) Oportunidad del día ({dia_code}) - Portafolio 60% a 79%")
    y -= 0.6 * cm
    c.setFont("Helvetica", 10)

    opp = g[(pd.to_numeric(g["PORTAFOLIO_PCT"], errors="coerce").notna()) &
            (g["PORTAFOLIO_PCT"] >= 60) & (g["PORTAFOLIO_PCT"] < 80)].copy()
    if "frecuencia" in opp.columns:
        od = opp[opp["frecuencia"].apply(lambda x: freq_has_day(x, dia_code))].copy()
        if not od.empty:
            opp = od
    opp = opp[opp["Vendedor_ID"].astype(str) != "1"].copy()

    if opp.empty:
        c.drawString(margin + 0.4 * cm, y, "Sin clientes en el rango 60–79% para el día.")
        y -= 0.6 * cm
    else:
        total_opp = int(opp["Cliente_ID"].nunique())
        c.drawString(margin + 0.4 * cm, y, f"Clientes únicos 60–79%: {total_opp}")
        y -= 0.5 * cm
        cnt = opp.groupby("Vendedor_ID")["Cliente_ID"].nunique().to_dict()
        for vend, val in sorted(cnt.items(), key=lambda kv: (-kv[1], str(kv[0]))):
            ensure_space(0.45 * cm)
            c.drawString(margin + 0.8 * cm, y, f"{vendor_display_name(vend)}: {int(val)}")
            y -= 0.42 * cm
        y -= 0.15 * cm

    # 6) Top 3 faltantes
    ensure_space(3.0 * cm)
    c.setFont("Helvetica-Bold", 11)
    c.drawString(margin, y, "6) Top 3 artículos más faltantes en portafolios TP")
    y -= 0.6 * cm
    c.setFont("Helvetica", 10)
    mode = matching_mode_label(g)
    top3 = top_missing_items_exact(g, top_n=3)
    if not top3:
        c.drawString(margin + 0.4 * cm, y, f"No disponible en esta corrida (modo de matching: {mode or 'SIN_DATO'}).")
        y -= 0.5 * cm
    else:
        for sku, clientes in top3:
            ensure_space(0.45 * cm)
            c.drawString(margin + 0.4 * cm, y, f"{sku} — falta en {clientes} clientes")
            y -= 0.42 * cm
    y -= 0.18 * cm

    # 7) Foco productos para la jornada
    ensure_space(3.0 * cm)
    c.setFont("Helvetica-Bold", 11)
    c.drawString(margin, y, f"7) Foco productos para {dia_code}")
    y -= 0.6 * cm
    c.setFont("Helvetica", 10)
    focus3 = top_focus_items_for_day(g, dia_code, top_n=3)
    if not focus3:
        c.drawString(margin + 0.4 * cm, y, f"No disponible en esta corrida (modo de matching: {mode or 'SIN_DATO'}).")
        y -= 0.5 * cm
    else:
        for sku, clientes in focus3:
            ensure_space(0.45 * cm)
            c.drawString(margin + 0.4 * cm, y, f"{sku} — foco en {clientes} clientes")
            y -= 0.42 * cm
    y -= 0.18 * cm

    # 8) Potencial
    ensure_space(3.5 * cm)
    c.setFont("Helvetica-Bold", 11)
    c.drawString(margin, y, "8) Potencial de facturación (sin IVA) - Planificación mañana (por UNIDAD)")
    y -= 0.55 * cm

    pot = compute_potential_60_79_for_tomorrow(g, dia_code, price_list)
    meta = price_list.get("meta", {}) or {}

    if not pot.get("has_price_data", False):
        c.setFont("Helvetica", 10)
        c.drawString(margin + 0.4 * cm, y, f"Lista precios no cargada. source={meta.get('source','')} path={meta.get('path','')}")
        y -= 0.55 * cm
    else:
        c.setFont("Helvetica", 10)
        c.drawString(margin + 0.4 * cm, y, f"Fuente: {pot.get('source','')} | Ítems: {pot.get('items',0)}")
        y -= 0.45 * cm
        c.drawString(margin + 0.4 * cm, y, f"Faltantes: {pot.get('falt',0)} | Valorizados: {pot.get('valued',0)} | Sin match: {pot.get('unknown',0)}")
        y -= 0.55 * cm
        c.setFont("Helvetica-Bold", 11)
        c.drawString(margin + 0.4 * cm, y, f"Potencial: {fmt_pesos(pot.get('pot_total',0.0))}")
        y -= 0.60 * cm

    finalize_page(c, W, H)
    c.save()
    return pdf_path





# ============================================================
# Export AppSheet CSVs
# ============================================================
def export_appsheet_csvs(g: pd.DataFrame, dia_code: str, base_out_dir: str,
                         pdf_vend_paths=None, pdf_ger_path: str = "",
                         generated_ts: str = ""):
    pdf_vend_paths = pdf_vend_paths or []
    generated_ts = generated_ts or datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    today_str = date.today().isoformat()

    appsheet_dir = os.path.join(base_out_dir, "APPSHEET")
    safe_mkdir(appsheet_dir)

    # ----------------------------
    # 1) pdf_index.csv
    # ----------------------------
    pdf_rows = []

    if pdf_ger_path and os.path.exists(pdf_ger_path):
        pdf_rows.append({
            "Fecha": today_str,
            "Dia": dia_code,
            "TipoPDF": "gerencia",
            "VendedorID": "",
            "VendedorNombre": "Gerencia",
            "ArchivoPDF": os.path.basename(pdf_ger_path),
            "RutaPDF": os.path.abspath(pdf_ger_path),
            "RutaPDF_Relativa": os.path.relpath(pdf_ger_path, start=base_out_dir),
            "Zona": dia_code,
            "Rol": "gerencia",
            "Visible": "SI",
            "GeneradoEn": generated_ts,
        })

    for p in pdf_vend_paths:
        if not p or not os.path.exists(p):
            continue
        m = re.search(r"Vendedor_(\d+)_" + re.escape(dia_code), os.path.basename(p))
        vend_id = m.group(1) if m else ""
        pdf_rows.append({
            "Fecha": today_str,
            "Dia": dia_code,
            "TipoPDF": "vendedor",
            "VendedorID": vend_id,
            "VendedorNombre": vendor_display_name(vend_id),
            "ArchivoPDF": os.path.basename(p),
            "RutaPDF": os.path.abspath(p),
            "RutaPDF_Relativa": os.path.relpath(p, start=base_out_dir),
            "Zona": dia_code,
            "Rol": "vendedor",
            "Visible": "SI",
            "GeneradoEn": generated_ts,
        })

    pd.DataFrame(pdf_rows).to_csv(os.path.join(appsheet_dir, "pdf_index.csv"), index=False, encoding="utf-8-sig")

    # ----------------------------
    # 2) clientes_oportunidad.csv
    # ----------------------------
    opp = g[(pd.to_numeric(g["PORTAFOLIO_PCT"], errors="coerce").notna()) &
            (g["PORTAFOLIO_PCT"] >= 60) & (g["PORTAFOLIO_PCT"] < 80)].copy()
    if "frecuencia" in opp.columns:
        od = opp[opp["frecuencia"].apply(lambda x: freq_has_day(x, dia_code))].copy()
        if not od.empty:
            opp = od
    opp = opp[opp["Vendedor_ID"].astype(str) != "1"].copy()

    vend_name_col = choose_vendor_name_column(opp) if not opp.empty else None

    opp_rows = []
    for _, r in opp.iterrows():
        vend_id = norm_vendor_id(r.get("Vendedor_ID", ""))
        vend_label = ""
        if vend_name_col and vend_name_col in opp.columns:
            vend_label = norm_text(r.get(vend_name_col, ""))
        vend_label = vendor_display_name(vend_id, vend_label)
        f80 = r.get("SKUs_faltan_para_80_ENT", [])
        if not isinstance(f80, list):
            f80 = parse_listish(f80)
        f100 = r.get("SKUs_faltan_para_100_ENT", [])
        if not isinstance(f100, list):
            f100 = parse_listish(f100)

        opp_rows.append({
            "Fecha": today_str,
            "Dia": dia_code,
            "VendedorID": vend_id,
            "VendedorNombre": vend_label,
            "ClienteID": int(r["Cliente_ID"]) if pd.notna(r.get("Cliente_ID")) else "",
            "RazonSocial": norm_text(r.get("Razon_Social", "")),
            "Localidad": norm_text(r.get("Localidad", "")),
            "Direccion": norm_text(r.get("Direccion", "")),
            "PortafolioPct": float(pd.to_numeric(r.get("PORTAFOLIO_PCT"), errors="coerce")) if pd.notna(pd.to_numeric(r.get("PORTAFOLIO_PCT"), errors="coerce")) else "",
            "Faltan80": int(pd.to_numeric(r.get("FALTAN_80_N", 0), errors="coerce") or 0),
            "Faltan100": int(pd.to_numeric(r.get("FALTAN_100_N", 0), errors="coerce") or 0),
            "SKUsFaltan80": " | ".join([norm_text(x) for x in f80 if norm_text(x)]),
            "SKUsFaltan100": " | ".join([norm_text(x) for x in f100 if norm_text(x)]),
            "TP_Sistema": "SI" if bool(r.get("TP_SISTEMA", False)) else "NO",
            "TP_Aplica": "SI" if bool(r.get("TP_ELIGIBLE", False)) else "NO",
        })

    pd.DataFrame(opp_rows).to_csv(os.path.join(appsheet_dir, "clientes_oportunidad.csv"), index=False, encoding="utf-8-sig")

    # ----------------------------
    # 3) foco_productos.csv
    # ----------------------------
    focus_rows = []

    # foco general del día
    for rank, item in enumerate(top_focus_items_for_day(g, dia_code, top_n=10), start=1):
        sku, clientes = item
        focus_rows.append({
            "Fecha": today_str,
            "Dia": dia_code,
            "Scope": "general",
            "VendedorID": "",
            "VendedorNombre": "General",
            "Rank": rank,
            "Articulo": sku,
            "CantClientes": int(clientes),
        })

    # foco por vendedor
    base = g.copy()
    base = base[(pd.to_numeric(base.get("PORTAFOLIO_PCT", np.nan), errors="coerce").notna()) &
                (base["PORTAFOLIO_PCT"] >= 60) & (base["PORTAFOLIO_PCT"] < 80)].copy()
    if "frecuencia" in base.columns:
        dayf = base[base["frecuencia"].apply(lambda x: freq_has_day(x, dia_code))].copy()
        if not dayf.empty:
            base = dayf
    base = base[base["Vendedor_ID"].astype(str) != "1"].copy()
    vend_name_col = choose_vendor_name_column(base) if not base.empty else None

    for vend_id, sub in base.groupby(base["Vendedor_ID"].map(norm_vendor_id)):
        counter = Counter()
        for _, r in sub.iterrows():
            skus = r.get("SKUs_faltan_para_80_ENT", [])
            if not isinstance(skus, list):
                skus = parse_listish(skus)
            seen = set()
            for sku in skus:
                k = norm_text(sku)
                if not k or k in seen:
                    continue
                counter[k] += 1
                seen.add(k)
        vend_label = ""
        if vend_name_col and vend_name_col in sub.columns and not sub.empty:
            vend_label = norm_text(sub.iloc[0].get(vend_name_col, ""))
        vend_label = vendor_display_name(vend_id, vend_label)
        for rank, (sku, clientes) in enumerate(counter.most_common(10), start=1):
            focus_rows.append({
                "Fecha": today_str,
                "Dia": dia_code,
                "Scope": "vendedor",
                "VendedorID": vend_id,
                "VendedorNombre": vend_label,
                "Rank": rank,
                "Articulo": sku,
                "CantClientes": int(clientes),
            })

    pd.DataFrame(focus_rows).to_csv(os.path.join(appsheet_dir, "foco_productos.csv"), index=False, encoding="utf-8-sig")
    print(f"[APPSHEET] CSVs -> {appsheet_dir}")

# ============================================================
# Main
# ============================================================
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--ventas", required=True)
    ap.add_argument("--gescom", required=True)
    ap.add_argument("--portafolio", required=False)
    ap.add_argument("--out", required=True)
    ap.add_argument("--export-pdfs", action="store_true")
    ap.add_argument("--export-gerencial", action="store_true")
    ap.add_argument("--dia-manana", default=None)
    ap.add_argument("--lista-precios", default=None)
    ap.add_argument("--objetivos-tp", default=None)
    args = ap.parse_args()

    dia = (args.dia_manana or day_code_for_tomorrow()).upper()
    if dia == "DO":
        dia = "LU"

    out_xlsx = args.out
    base_out_dir = os.path.dirname(out_xlsx) if os.path.dirname(out_xlsx) else "."
    pdf_vend_dir = os.path.join(base_out_dir, "PDF_VENDEDORES", date.today().isoformat())
    pdf_ger_dir = os.path.join(base_out_dir, "PDF_GERENCIAL")

    hist_dir = os.path.join(base_out_dir, "cache_tp_portafolio")
    hist_path = os.path.join(hist_dir, "hist_tp_runs.json")
    debug_path = os.path.join(base_out_dir, "debug_tp_compra.txt")

    ventas_path = resolve_input_file(args.ventas, [r"venta.*\.(csv|xlsx|xls)$", r".*ventas.*\.(csv|xlsx|xls)$"])
    gescom_path = resolve_input_file(args.gescom, [r"gescom.*\.(xlsx|xls)$"])
    if not ventas_path or not os.path.exists(ventas_path):
        raise FileNotFoundError(f"No se encontró archivo de ventas. Recibido: {args.ventas}")
    if not gescom_path or not os.path.exists(gescom_path):
        raise FileNotFoundError(f"No se encontró archivo Gescom. Recibido: {args.gescom}")

    # ventas + gescom
    sales = load_sales_clean(ventas_path)
    ges = load_gescom(gescom_path)
    ges = attach_localidad_from_sales(ges, sales)
    ges = attach_sales_metrics_from_sales(ges, sales)
    try:
        print(f"[VENTAS] archivo={ventas_path} columna_importe={ges['_VENTA_MONTO_COL'].iloc[0] if '_VENTA_MONTO_COL' in ges.columns and not ges.empty else ''} fact_total={ges['FACTURACION_CLIENTE'].sum():.2f}")
    except Exception:
        pass
    g = compute_tp_flags(ges)

    portafolio_path = resolve_input_file(args.portafolio, [r"portafolio.*\.pptx$", r".*infaltable.*\.pptx$", r".*portfolio.*\.pptx$"])
    if portafolio_path and os.path.exists(portafolio_path):
        print(f"[PORTAFOLIO] OK -> {portafolio_path}")
    else:
        print(f"[PORTAFOLIO] No encontrado. Recibido: {args.portafolio}")

    # faltantes (y debug)
    g = attach_portafolio_missing_to_g(g, sales, portafolio_path, debug_path)

    # historial deltas
    runs = load_hist_runs(hist_path)
    runs, snap, delta_prev_run, delta_prev_day = update_hist_and_get_deltas(g, runs, hist_path)

    # objetivos TP
    objetivos_path = args.objetivos_tp
    if not objetivos_path:
        c1 = os.path.join(script_dir(), "inputs", "objetivo_vendedor_tp.xlsx")
        c2 = os.path.join(script_dir(), "inputs", "objetivo_vendedor_tp.xls")
        if os.path.exists(c1):
            objetivos_path = c1
        elif os.path.exists(c2):
            objetivos_path = c2
    else:
        if not os.path.isabs(objetivos_path):
            objetivos_path = os.path.join(script_dir(), objetivos_path)

    prev_run_snap_for_obj = runs[-2]["snapshot"] if len(runs) >= 2 and "snapshot" in runs[-2] else None
    prev_day_snap_for_obj = None
    for r in reversed(runs[:-1]):
        ts = str(r.get("ts",""))
        if ts[:10] < date.today().isoformat() and "snapshot" in r:
            prev_day_snap_for_obj = r["snapshot"]
            break
    obj_bundle = load_objectives_tp(objetivos_path, g, prev_day_snap_for_obj) if objetivos_path else {"loaded": False, "error": "Sin archivo de objetivos."}
    if obj_bundle.get("loaded"):
        print(f"[OBJETIVOS_TP] OK -> {objetivos_path}")
    else:
        print(f"[OBJETIVOS_TP] {obj_bundle.get('error','No cargado')}")

    # lista precios
    price_path = args.lista_precios
    if not price_path:
        c1 = os.path.join(script_dir(), "inputs", "lista_precios.xlsx")
        c2 = os.path.join(script_dir(), "inputs", "lista_precios.csv")
        c3 = os.path.join(script_dir(), "inputs", "lista_precios.xls")
        if os.path.exists(c1):
            price_path = c1
        elif os.path.exists(c2):
            price_path = c2
        elif os.path.exists(c3):
            price_path = c3
    else:
        if not os.path.isabs(price_path):
            price_path = os.path.join(script_dir(), price_path)

    price_list = load_price_list_from_table(price_path) if price_path else {"by_desc": {}, "meta": {"source": "SinListaPreciosTabla", "items": 0, "path": ""}}
    meta = price_list.get("meta", {}) or {}
    print(f"[PRECIOS_TABLA] source={meta.get('source','')} items={meta.get('items',0)} path={meta.get('path','')}")
    print(f"[DEBUG] -> {debug_path}")
    print(f"[TP_RUNNER] version={SCRIPT_VERSION} script={os.path.abspath(__file__)}")

    # exports
    created = []
    pdf_ger = ""
    if args.export_pdfs:
        created = export_pdf_vendedores(g, dia, pdf_vend_dir, obj_bundle=obj_bundle)
        print(f"PDFs vendedores: {len(created)} -> {pdf_vend_dir}")

    if args.export_gerencial:
        pdf_ger = export_pdf_gerencial(g, snap, delta_prev_run, delta_prev_day, dia, pdf_ger_dir, price_list, obj_bundle=obj_bundle)
        print(f"Gerencial -> {pdf_ger}")

    export_appsheet_csvs(
        g=g,
        dia_code=dia,
        base_out_dir=base_out_dir,
        pdf_vend_paths=created,
        pdf_ger_path=pdf_ger,
        generated_ts=datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
    )

    # xlsx control
    safe_mkdir(base_out_dir)
    try:
        with pd.ExcelWriter(out_xlsx, engine="openpyxl") as w:
            g.to_excel(w, index=False, sheet_name="STATUS_CLIENTE")
            if obj_bundle.get("loaded"):
                obj_bundle.get("tax_df", pd.DataFrame()).to_excel(w, index=False, sheet_name="OBJ_TAXONOMIA")
                obj_bundle.get("vend_df", pd.DataFrame()).to_excel(w, index=False, sheet_name="OBJ_VENDEDOR")
        print("OK ->", out_xlsx)
    except Exception as e:
        print("WARN xlsx:", str(e))


if __name__ == "__main__":
    main()